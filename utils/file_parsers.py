import os
import logging
from typing import List, Optional, Dict, Any, Tuple
from abc import ABC, abstractmethod

# 設置日誌
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 自定義過濾器，過濾不重要的警告消息
class EnhancedWarningFilter(logging.Filter):
    def __init__(self):
        super().__init__()
        self.warned_about = set()
        self.error_counts = {}  # 添加錯誤計數器
        
    def filter(self, record):
        # 處理錯誤級別訊息
        if record.levelno == logging.ERROR:
            message = record.getMessage().lower()
            
            # 對於預期內的降級行為，將其降級為 INFO 級別
            expected_fallbacks = [
                "初始化 word 應用失敗",
                "wdalertsnone",
                "pagesetup",
                "不是標準ole2文件",
                "com 接口",
                "module 'olefile' has no attribute",  # 添加 olefile 相關錯誤
                "解析 visio 文件",                    # 添加 Visio 解析相關錯誤
                "open() got an unexpected keyword"   # 添加關鍵字參數錯誤
            ]
            
            for fallback_text in expected_fallbacks:
                if fallback_text in message:
                    # 將錯誤降級為 INFO
                    record.levelno = logging.INFO
                    record.levelname = "INFO"
                    
                    # 修改訊息使其更友好
                    if "初始化 word 應用失敗" in message:
                        record.msg = "COM 接口不可用，將嘗試後備方法: %s"
                    elif "module 'olefile' has no attribute" in message:
                        record.msg = "Visio 解析使用後備方法: %s"
                    
                    # 計數相同錯誤
                    key = fallback_text
                    self.error_counts[key] = self.error_counts.get(key, 0) + 1
                    
                    # 如果同一錯誤出現過多次，不再顯示
                    if self.error_counts[key] > 2:
                        return False
                    break
            
        # 過濾掉常見的非關鍵警告
        if record.levelno == logging.WARNING:
            message = record.getMessage().lower()
            
            # 過濾 openpyxl 的標頭和頁腳解析警告
            if "header or footer" in message:
                return False
                
            # 以下是只需要顯示一次的警告
            one_time_warnings = [
                "未安裝 textract 庫",
                "系統中未找到 antiword 工具",
                "找不到 libreoffice",
                "非 windows 系統",
                "不是標準 ole2 格式"
            ]
            
            # 檢查是否為已經顯示過的一次性警告
            for warning_text in one_time_warnings:
                if warning_text in message:
                    # 如果已經警告過，則不再顯示
                    warning_key = warning_text
                    if warning_key in self.warned_about:
                        return False
                    self.warned_about.add(warning_key)
                    break
                    
        return True

# 應用增強型過濾器到日誌器
logger.addFilter(EnhancedWarningFilter())

# 導入基本依賴
import fitz  # PyMuPDF
from docx import Document
import openpyxl
import csv

# 在導入部分添加新的依賴
import tempfile
import subprocess
import io
from pathlib import Path
import warnings
import platform
import shutil

# 嘗試初始化 COM 組件（用於 Windows 環境）
try:
    if platform.system() == 'Windows':
        import pythoncom
        pythoncom.CoInitialize()
        logger.info("COM 組件已成功初始化")
except ImportError:
    logger.warning("無法導入 pythoncom 模塊，某些 Windows 文件解析功能可能受限")
except Exception as e:
    logger.warning(f"初始化 COM 組件時出錯: {str(e)}")

# 設置警告過濾
warnings.filterwarnings("ignore", category=DeprecationWarning)

# 移除全局 OCR 導入
TESSERACT_AVAILABLE = False

# 解析方法統計和可用性追蹤類
class ParserMethodStats:
    """追蹤文件解析方法的可用性和成功率"""
    
    def __init__(self):
        self.available_methods = {}
        self.success_stats = {}
        self.method_checked = False
    
    def is_method_available(self, method_name):
        """檢查解析方法是否可用"""
        if not self.method_checked:
            self._check_all_methods()
        return self.available_methods.get(method_name, False)
    
    def record_success(self, method_name):
        """記錄方法成功解析"""
        if method_name not in self.success_stats:
            self.success_stats[method_name] = {'success': 0, 'attempts': 0}
        self.success_stats[method_name]['success'] += 1
        self.success_stats[method_name]['attempts'] += 1
    
    def record_failure(self, method_name):
        """記錄方法解析失敗"""
        if method_name not in self.success_stats:
            self.success_stats[method_name] = {'success': 0, 'attempts': 0}
        self.success_stats[method_name]['attempts'] += 1
    
    def get_success_rate(self, method_name):
        """獲取方法的成功率"""
        if method_name not in self.success_stats:
            return 0.0
        stats = self.success_stats[method_name]
        if stats['attempts'] == 0:
            return 0.0
        return stats['success'] / stats['attempts']
    
    def _check_all_methods(self):
        """檢查所有解析方法的可用性"""
        # 檢查 COM 可用性
        self.available_methods['_parse_with_com_win32'] = platform.system() == 'Windows'
        
        # 檢查 LibreOffice 可用性
        libreoffice_paths = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            '/usr/bin/libreoffice',
            '/usr/bin/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        ]
        self.available_methods['_parse_with_libreoffice'] = any(os.path.exists(path) for path in libreoffice_paths)
        
        # 檢查 textract 可用性
        try:
            import textract
            self.available_methods['_parse_with_textract'] = True
        except ImportError:
            self.available_methods['_parse_with_textract'] = False
        
        # 檢查 antiword 可用性
        try:
            result = subprocess.run(
                ["where", "antiword"] if platform.system() == "Windows" else ["which", "antiword"], 
                capture_output=True, 
                text=True, 
                check=False
            )
            self.available_methods['_parse_with_antiword'] = result.returncode == 0
        except Exception:
            self.available_methods['_parse_with_antiword'] = False
        
        # 二進制讀取方法總是可用的
        self.available_methods['_parse_with_binary_read'] = True
        self.available_methods['_binary_force_extract'] = True
        
        # 記錄可用方法
        available = [name for name, available in self.available_methods.items() if available]
        logger.info(f"可用的文件解析方法: {', '.join(available)}")
        
        # 記錄不可用方法和安裝建議
        unavailable = [name for name, available in self.available_methods.items() if not available]
        if unavailable:
            install_suggestions = []
            if '_parse_with_textract' in unavailable:
                install_suggestions.append("pip install textract - 提供更好的文本提取能力")
            if '_parse_with_libreoffice' in unavailable:
                install_suggestions.append("安裝 LibreOffice - 提高 Office 文件的解析質量")
            if '_parse_with_antiword' in unavailable:
                install_suggestions.append("安裝 antiword 工具 - 改進舊版 Word 文件解析")
            
            if install_suggestions:
                logger.info(f"文件解析工具安裝建議:\n" + "\n".join(install_suggestions))
        
        self.method_checked = True
        return self.available_methods

# 創建解析方法統計實例
parser_stats = ParserMethodStats()

class FileParser(ABC):
    """
    文件解析器的抽象基類。
    所有具體的解析器類都必須繼承此類並實現 parse 方法。
    """
    
    def __init__(self) -> None:
        self.file_path: Optional[str] = None

    @abstractmethod
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析文件並返回文本塊和元數據的列表。

        Args:
            file_path (str): 要解析的文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        pass

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        提取文件的基本元數據。

        Args:
            file_path (str): 文件路徑

        Returns:
            Dict[str, Any]: 包含文件元數據的字典
        """
        try:
            stats = os.stat(file_path)
            return {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_size": stats.st_size,
                "last_modified": stats.st_mtime
            }
        except Exception as e:
            logger.error(f"提取元數據時出錯 {file_path}: {str(e)}")
            return {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "error": str(e)
            }

    @staticmethod
    def get_parser_for_file(file_path: str) -> Optional['FileParser']:
        """
        根據文件擴展名返回適當的解析器實例。

        Args:
            file_path (str): 文件路徑

        Returns:
            Optional[FileParser]: 解析器實例，如果不支持該文件類型則返回 None
        """
        extension = os.path.splitext(file_path)[1].lower()
        parser_map = {
            '.pdf': PDFParser,
            '.docx': DocxParser,
            '.doc': DocParser,
            '.xlsx': ExcelParser,
            '.xls': ExcelOldParser,
            '.txt': TextParser,
            '.md': TextParser,
            '.pptx': PPTXParser,
            '.ppt': PPTParser,
            '.csv': CSVParser,
            '.vsdx': VisioParser,
            '.vsd': VisioParser
        }
        
        parser_class = parser_map.get(extension)
        if parser_class:
            try:
                parser = parser_class()
                parser.file_path = file_path
                return parser
            except Exception as e:
                logger.error(f"創建解析器實例時出錯 {file_path}: {str(e)}")
                return None
        else:
            logger.warning(f"不支持的文件類型: {extension}")
            return None

    def safe_parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        安全地解析文件，捕獲所有異常並返回友好的錯誤消息

        Args:
            file_path: 文件路徑
            
        Returns:
            包含 (文本, 元數據) 元組的列表
        """
        try:
            return self.parse(file_path)
        except Exception as e:
            logger.error(f"解析文件時出錯 ({file_path}): {str(e)}")
            import traceback
            error_detail = traceback.format_exc()
            logger.debug(f"詳細錯誤: {error_detail}")
            
            metadata = self.extract_metadata(file_path)
            metadata["error"] = str(e)
            metadata["error_detail"] = error_detail
            
            return [("無法解析此文件，可能是格式不兼容或文件損壞", metadata)]


class PDFParser(FileParser):
    """PDF 文件解析器"""

    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析 PDF 文件並返回文本內容和元數據。

        Args:
            file_path (str): PDF 文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        results = []
        metadata = self.extract_metadata(file_path)
        metadata["document_type"] = "PDF"

        try:
            # 設置 PyMuPDF 的選項，忽略一些錯誤
            # 這有助於解決 Pdfpwwin32 相關的錯誤
            fitz.TOOLS.mupdf_display_errors(False)
            
            with fitz.open(file_path) as doc:
                # 收集所有頁面的文本
                all_text = []
                
                for page_num, page in enumerate(doc, 1):
                    try:
                        # 嘗試使用不同的文本提取方法
                        text = page.get_text()
                        
                        if not text.strip():
                            # 嘗試用不同的方法提取，如 "blocks" 模式
                            try:
                                blocks = page.get_text("blocks")
                                if blocks:
                                    text = "\n".join([b[4] for b in blocks])
                            except Exception:
                                pass
                        
                        if text.strip():
                            all_text.append(text)
                            
                            # 每 5 頁或文件結束時創建一個塊
                            if len(all_text) >= 5 or page_num == len(doc):
                                combined_text = "\n\n".join(all_text)
                                page_metadata = metadata.copy()
                                page_metadata.update({
                                    "page_range": f"{page_num-len(all_text)+1}-{page_num}",
                                    "total_pages": len(doc)
                                })
                                results.append((combined_text, page_metadata))
                                all_text = []
                                
                    except Exception as e:
                        # 特別處理 Pdfpwwin32 錯誤
                        error_msg = str(e)
                        if "Pdfpwwin32" in error_msg:
                            logger.warning(f"Pdfpwwin32爬取文件失敗: {file_path} - 這通常是由於PDF權限限制導致")
                        else:
                            logger.warning(f"處理 PDF 頁面 {page_num} 時出錯: {str(e)}")
                        continue
                
                # 確保最後的文本塊被處理
                if all_text:
                    combined_text = "\n\n".join(all_text)
                    page_metadata = metadata.copy()
                    page_metadata.update({
                        "page_range": f"remaining",
                        "total_pages": len(doc)
                    })
                    results.append((combined_text, page_metadata))

        except Exception as e:
            error_msg = str(e)
            if "Pdfpwwin32" in error_msg or "CoInitialize" in error_msg:
                logger.error(f"PDF權限問題或COM初始化錯誤: {file_path} - {error_msg}")
                metadata["error"] = "PDF權限問題或COM初始化錯誤"
            else:
                logger.error(f"解析 PDF 文件 {file_path} 時出錯: {error_msg}")
                metadata["error"] = error_msg
            
            results.append(("PDF 解析出錯，可能是文件格式問題或權限限制", metadata))

        return results if results else [("PDF 文件為空或無法提取文本", metadata)]


class DocxParser(FileParser):
    """Word 文件解析器"""

    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析 Word 文件並返回文本內容和元數據。

        Args:
            file_path (str): Word 文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        results = []
        metadata = self.extract_metadata(file_path)
        metadata["document_type"] = "DOCX"

        try:
            doc = Document(file_path)
            text_blocks = []
            paragraph_count = 0

            # 提取段落文本
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    text_blocks.append(text)
                    paragraph_count += 1

            # 提取表格文本
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_values = [cell.text.strip() for cell in row.cells]
                    if any(row_values):
                        import re
                        row_text = " | ".join(row_values)
                        cleaned_row_text = re.sub(r'(\s*\|\s*){2,}', ' | ', row_text).strip(' |')
                        if cleaned_row_text:
                            table_text.append(cleaned_row_text)
                if table_text:
                    text_blocks.append("\n".join(table_text))

            if text_blocks:
                metadata.update({
                    "paragraph_count": paragraph_count
                })
                results.append(("\n\n".join(text_blocks), metadata))

        except Exception as e:
            logger.error(f"解析 Word 文件 {file_path} 時出錯: {str(e)}")
            metadata["error"] = str(e)
            results.append(("Word 文件解析出錯，可能是文件格式問題或權限限制", metadata))

        return results if results else [("Word 文件為空或無法提取文本", metadata)]


class ExcelParser(FileParser):
    """Excel 文件解析器"""

    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析 Excel 文件並返回文本內容和元數據。

        Args:
            file_path (str): Excel 文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        results = []
        metadata = self.extract_metadata(file_path)
        metadata["document_type"] = "XLSX"

        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            sheet_processed = False

            for sheet_name in workbook.sheetnames:
                try:
                    sheet = workbook[sheet_name]
                    sheet_text = []
                    
                    # 提取表格數據，設置行數限制以防止超大表格
                    max_rows = 5000
                    row_count = 0
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_count += 1
                        if row_count > max_rows:
                            sheet_text.append("... (表格過大，僅顯示前5000行)")
                            break
                        
                        try:
                            row_values = [str(cell).strip() if cell is not None else "" for cell in row]
                            if any(row_values):
                                import re
                                row_text = " | ".join(row_values)
                                # 將多個分隔符壓縮為一個，並清理開頭和結尾的分隔符
                                cleaned_row_text = re.sub(r'(\s*\|\s*){2,}', ' | ', row_text).strip(' |')
                                # 只有在清理後仍有內容時才添加
                                if cleaned_row_text:
                                    sheet_text.append(cleaned_row_text)
                        except Exception as row_error:
                            logger.warning(f"處理 Excel 行 {row_count} 時出錯: {str(row_error)}")
                            continue
                        
                    if sheet_text:
                        sheet_content = "\n".join(sheet_text)
                        sheet_metadata = metadata.copy()
                        sheet_metadata.update({
                            "sheet_name": sheet_name,
                            "row_count": row_count
                        })
                        results.append((sheet_content, sheet_metadata))
                        sheet_processed = True

                except Exception as sheet_error:
                    logger.warning(f"處理 Excel 表格 {sheet_name} 時出錯: {str(sheet_error)}")
                    continue

            workbook.close()

        except Exception as e:
            logger.error(f"解析 Excel 文件 {file_path} 時出錯: {str(e)}")
            metadata["error"] = str(e)
            results.append(("Excel 文件解析出錯，可能是文件格式問題或權限限制", metadata))

        return results if results else [("Excel 文件為空或無法提取文本", metadata)]


class TextParser(FileParser):
    """文本文件解析器 (TXT, MD等) - 支持多種編碼自動檢測"""
    
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析文本文件並返回文本內容和元數據。

        Args:
            file_path (str): 文本文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        results = []
        metadata = self.extract_metadata(file_path)
        metadata["document_type"] = os.path.splitext(file_path)[1][1:].upper() or "TXT"

        # 嘗試多種編碼
        encodings = ['utf-8', 'big5', 'gb18030', 'gb2312', 'latin-1', 'cp1252', 'utf-16']
        content = None
        used_encoding = None
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding, errors='strict') as f:
                    content = f.read()
                    used_encoding = encoding
                    break
            except (UnicodeDecodeError, UnicodeError):
                continue
            except Exception as e:
                logger.warning(f"使用編碼 {encoding} 讀取文件時出錯: {str(e)}")
                continue
        
        # 如果所有編碼都失敗，使用二進制讀取並嘗試修復
        if content is None:
            try:
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                    # 嘗試使用chardet檢測編碼
                    try:
                        import chardet
                        detected = chardet.detect(raw_data)
                        if detected['encoding'] and detected['confidence'] > 0.7:
                            content = raw_data.decode(detected['encoding'], errors='replace')
                            used_encoding = detected['encoding']
                            logger.info(f"使用chardet檢測到編碼: {detected['encoding']} (置信度: {detected['confidence']:.2f})")
                    except ImportError:
                        pass
                    
                    # 如果chardet也失敗，使用utf-8並替換錯誤字符
                    if content is None:
                        content = raw_data.decode('utf-8', errors='replace')
                        used_encoding = 'utf-8 (with errors replaced)'
                        
            except Exception as e:
                logger.error(f"二進制讀取文件 {file_path} 時出錯: {str(e)}")
                metadata["error"] = str(e)
                results.append(("文本文件解析出錯，可能是編碼問題或權限限制", metadata))
                return results

        if content and content.strip():
            metadata["line_count"] = len(content.splitlines())
            metadata["encoding"] = used_encoding
            metadata["file_size_chars"] = len(content)
            
            # 檢測並清理可能的亂碼
            cleaned_content = self._clean_garbled_text(content)
            results.append((cleaned_content, metadata))
        else:
            results.append(("文本文件為空或無法提取文本", metadata))

        return results
    
    def _clean_garbled_text(self, text: str) -> str:
        """
        清理可能的亂碼文本
        
        Args:
            text: 原始文本
            
        Returns:
            清理後的文本
        """
        if not text:
            return text
        
        # 移除常見的亂碼字符
        garbled_chars = ['\ufffd', '\x00', '\x01', '\x02', '\x03', '\x04', '\x05', '\x06', '\x07', '\x08']
        cleaned_text = text
        
        for char in garbled_chars:
            cleaned_text = cleaned_text.replace(char, '')
        
        # 移除過多的空白字符
        import re
        cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)  # 最多保留兩個連續換行
        cleaned_text = re.sub(r' {3,}', '  ', cleaned_text)     # 最多保留兩個連續空格
        
        return cleaned_text.strip()


class CSVParser(FileParser):
    """CSV 文件解析器"""
    
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析 CSV 文件並返回文本內容和元數據。

        Args:
            file_path (str): CSV 文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        results = []
        metadata = self.extract_metadata(file_path)
        metadata["document_type"] = "CSV"
        
        # 嘗試不同的編碼
        encodings = ['utf-8', 'big5', 'gb18030', 'latin-1']
        success = False
        
        for encoding in encodings:
            try:
                rows = []
                row_count = 0
                
                with open(file_path, 'r', encoding=encoding, newline='') as f:
                    reader = csv.reader(f)
                    for row in reader:
                        row_count += 1
                        if row_count > 5000:  # 限制行數
                            rows.append("... (表格過大，僅顯示前5000行)")
                            break
                        
                        row_values = [str(cell).strip() for cell in row]
                        if any(row_values):
                            import re
                            row_text = " | ".join(row_values)
                            cleaned_row_text = re.sub(r'(\s*\|\s*){2,}', ' | ', row_text).strip(' |')
                            if cleaned_row_text:
                                rows.append(cleaned_row_text)
                
                if rows:
                    content = "\n".join(rows)
                    metadata.update({
                        "encoding": encoding,
                        "row_count": row_count
                    })
                    results.append((content, metadata))
                    success = True
                    break
                    
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.warning(f"使用編碼 {encoding} 解析 CSV 文件時出錯: {str(e)}")
                continue
        
        if not success:
            metadata["error"] = "無法使用支持的編碼解析文件"
            results.append(("CSV 文件解析出錯，可能是編碼問題或格式錯誤", metadata))
            
        return results if results else [("CSV 文件為空或無法提取文本", metadata)]


class PPTXParser(FileParser):
    """PowerPoint 文件解析器 (PPTX)"""
    
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析 PowerPoint 文件並返回文本內容和元數據。

        Args:
            file_path (str): PowerPoint 文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        results = []
        metadata = self.extract_metadata(file_path)
        metadata["document_type"] = "PPTX"
        
        try:
            from pptx import Presentation

            # 打開演示文稿
            presentation = Presentation(file_path)
            
            # 用於存儲提取的文本
            all_text = []
            
            # 添加標題頁元數據
            if presentation.core_properties.title:
                metadata["title"] = presentation.core_properties.title
            if presentation.core_properties.author:
                metadata["author"] = presentation.core_properties.author
            if presentation.core_properties.created:
                metadata["created"] = str(presentation.core_properties.created)
            
            # 提取幻燈片數量
            metadata["slide_count"] = len(presentation.slides)
            
            # 遍歷所有幻燈片
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                slide_text.append(f"===== 幻燈片 {i+1} =====")
                
                # 提取幻燈片標題 (如果有)
                if slide.shapes.title and slide.shapes.title.text:
                    slide_text.append(f"標題: {slide.shapes.title.text}")
                
                # 從所有形狀中提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # 跳過已經處理過的標題
                        if shape == slide.shapes.title:
                            continue
                        slide_text.append(shape.text)
                
                # 如果幻燈片有提取到內容，加入到總文本中
                if len(slide_text) > 1:  # 超過了僅有的幻燈片標題行
                    all_text.append("\n".join(slide_text))
            
            # 將所有內容合併為一個字符串，並添加到結果中
            if all_text:
                results.append(("\n\n".join(all_text), metadata))
            else:
                results.append(("PowerPoint 文件未包含可提取的文本", metadata))
                
        except ImportError:
            error_msg = "未能導入 python-pptx 庫，無法解析 PPTX 文件"
            logger.error(error_msg)
            metadata["error"] = error_msg
            results.append((error_msg, metadata))
        except Exception as e:
            error_msg = f"解析 PowerPoint 文件 {file_path} 時出錯: {str(e)}"
            logger.error(error_msg)
            metadata["error"] = str(e)
            results.append((error_msg, metadata))
            
        return results


class VisioParser(FileParser):
    """Visio 文件解析器"""
    
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析 Visio 文件並返回文本內容和元數據。

        Args:
            file_path (str): Visio 文件路徑

        Returns:
            List[Tuple[str, Dict[str, Any]]]: 包含 (文本, 元數據) 元組的列表
        """
        results = []
        metadata = self.extract_metadata(file_path)
        metadata["document_type"] = "VSD" if file_path.lower().endswith(".vsd") else "VSDX"
        
        try:
            # VSDX 文件實際上是 ZIP 文件，可以使用 zipfile 庫處理
            if file_path.lower().endswith(".vsdx"):
                import zipfile
                import xml.etree.ElementTree as ET
                
                with zipfile.ZipFile(file_path, 'r') as vsdx:
                    # 提取文件列表
                    metadata["file_entries"] = len(vsdx.namelist())
                    
                    # 提取文本內容
                    text_content = []
                    
                    # 嘗試從 document.xml 提取信息
                    try:
                        with vsdx.open('docProps/core.xml') as f:
                            core_xml = f.read()
                            root = ET.fromstring(core_xml)
                            
                            # 提取元數據
                            ns = {'cp': 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
                                 'dc': 'http://purl.org/dc/elements/1.1/'}
                            
                            # 提取標題
                            title_elem = root.find('.//dc:title', ns)
                            if title_elem is not None and title_elem.text:
                                metadata["title"] = title_elem.text
                                text_content.append(f"標題: {title_elem.text}")
                            
                            # 提取作者
                            creator_elem = root.find('.//dc:creator', ns)
                            if creator_elem is not None and creator_elem.text:
                                metadata["author"] = creator_elem.text
                                text_content.append(f"作者: {creator_elem.text}")
                            
                            # 提取說明
                            desc_elem = root.find('.//dc:description', ns)
                            if desc_elem is not None and desc_elem.text:
                                metadata["description"] = desc_elem.text
                                text_content.append(f"說明: {desc_elem.text}")
                    except Exception as xml_error:
                        logger.warning(f"解析 VSDX 元數據時出錯: {str(xml_error)}")
                    
                    # 查找 pages
                    pages = [f for f in vsdx.namelist() if f.startswith('visio/pages/page') and f.endswith('.xml')]
                    metadata["page_count"] = len(pages)
                    
                    # 從每個頁面提取文本
                    for page_file in pages:
                        try:
                            with vsdx.open(page_file) as f:
                                page_xml = f.read()
                                page_root = ET.fromstring(page_xml)
                                
                                # 提取頁面名稱
                                page_name = page_file.split('/')[-1].replace('.xml', '')
                                text_content.append(f"\n===== {page_name} =====")
                                
                                # 提取文本元素
                                ns = {'vis': 'http://schemas.microsoft.com/office/visio/2012/main'}
                                texts = page_root.findall('.//vis:Text', ns)
                                
                                for text_elem in texts:
                                    if text_elem.text and text_elem.text.strip():
                                        text_content.append(text_elem.text.strip())
                        except Exception as page_error:
                            logger.warning(f"解析 VSDX 頁面 {page_file} 時出錯: {str(page_error)}")
                    
                    # 如果沒有找到文本，嘗試使用常規方法搜索可能包含文本的文件
                    if len(text_content) <= 1:
                        text_files = [f for f in vsdx.namelist() if f.endswith('.txt') or f.endswith('.xml')]
                        for text_file in text_files[:10]:  # 限制處理文件數量
                            try:
                                with vsdx.open(text_file) as f:
                                    content = f.read().decode('utf-8', errors='ignore')
                                    if len(content) > 10:  # 忽略非常短的內容
                                        text_content.append(f"\n--- {text_file} ---")
                                        text_content.append(content[:1000])  # 限制每個文件的文本長度
                            except Exception:
                                pass
                
                # 如果有提取到內容，加入到結果中
                if text_content:
                    results.append(("\n".join(text_content), metadata))
                else:
                    results.append(("無法從 VSDX 文件提取文本內容", metadata))
            
            # VSD 文件使用 OLE 格式，需要 olefile 庫
            elif file_path.lower().endswith(".vsd"):
                import olefile
                
                if not olefile.isOleFile(file_path):
                    results.append(("文件不是有效的 OLE 格式的 Visio 文件", metadata))
                    return results
                
                text_content = []  # 在 try 塊外定義
                
                try:
                    # 修正函數名稱：OleFile -> OleFileIO
                    with olefile.OleFileIO(file_path) as ole:
                        # 獲取 OLE 流列表
                        streams = ole.listdir()
                        metadata["stream_count"] = len(streams)
                        
                        # 嘗試從 SummaryInformation 和 DocumentSummaryInformation 流中獲取元數據
                        for summary_stream in ['\\x05SummaryInformation', '\\x05DocumentSummaryInformation']:
                            if ole.exists(summary_stream):
                                try:
                                    # 獲取摘要信息
                                    summary = ole.getproperties(summary_stream)
                                    
                                    # 添加可能存在的元數據
                                    if 1 in summary and summary[1]:  # 標題
                                        metadata["title"] = summary[1]
                                        text_content.append(f"標題: {summary[1]}")
                                    if 2 in summary and summary[2]:  # 主題
                                        metadata["subject"] = summary[2]
                                        text_content.append(f"主題: {summary[2]}")
                                    if 3 in summary and summary[3]:  # 作者
                                        metadata["author"] = summary[3]
                                        text_content.append(f"作者: {summary[3]}")
                                    if 4 in summary and summary[4]:  # 關鍵字
                                        metadata["keywords"] = summary[4]
                                        text_content.append(f"關鍵字: {summary[4]}")
                                    if 5 in summary and summary[5]:  # 註釋
                                        metadata["comments"] = summary[5]
                                        text_content.append(f"註釋: {summary[5]}")
                                except Exception as sum_error:
                                    logger.warning(f"讀取 VSD 摘要信息時出錯: {str(sum_error)}")
                        
                        # 嘗試從其他流中提取文本
                        # Visio 文件中的文本可能存在於多個流中，我們嘗試幾個可能的位置
                        text_streams = []
                        
                        # 查找可能包含文本的流
                        for stream in streams:
                            stream_name = "/".join(stream)
                            if ("Text" in stream_name or 
                                "Datas" in stream_name or 
                                "Contents" in stream_name or
                                "VisioDocument" in stream_name):
                                text_streams.append(stream)
                        
                        # 從每個可能的流中提取文本
                        for stream in text_streams:
                            try:
                                with ole.openstream(stream) as s:
                                    data = s.read()
                                    # 嘗試以不同編碼解析文本
                                    for encoding in ['utf-8', 'utf-16-le', 'cp1252', 'gbk']:
                                        try:
                                            decoded = data.decode(encoding, errors='ignore')
                                            # 只保留可打印的 ASCII 字符和基本標點
                                            import re
                                            filtered = re.sub(r'[^\x20-\x7E\u4e00-\u9fff\u3040-\u30ff\u0100-\u017f\s.,!?;:-]', ' ', decoded)
                                            filtered = re.sub(r'\s+', ' ', filtered).strip()
                                            
                                            # 如果解析出有意義的文本，添加到結果中
                                            if len(filtered) > 20:  # 忽略太短的文本
                                                text_content.append(f"\n--- {'/'.join(stream)} ---")
                                                # 限制每個流的文本長度
                                                text_content.append(filtered[:2000])
                                                break
                                        except:
                                            continue
                            except Exception as stream_error:
                                logger.debug(f"讀取流 {'/'.join(stream)} 時出錯: {str(stream_error)}")
                except Exception as ole_error:
                    logger.warning(f"使用 OleFileIO 解析 VSD 文件時出錯: {str(ole_error)}")
                    # 即使失敗也繼續，我們可能已經提取了一些文本
                
                # 如果有提取到內容，加入到結果中 (現在在 try 區塊外)
                if len(text_content) > 1:  # 不僅僅是標題
                    results.append(("\n".join(text_content), metadata))
                else:
                    results.append(("無法從 VSD 文件提取文本內容", metadata))
            
            else:
                results.append((f"不支持的 Visio 文件格式: {file_path}", metadata))
            
        except ImportError as ie:
            error_msg = f"未能導入所需庫進行 Visio 解析: {str(ie)}"
            logger.warning(error_msg)  # 將 ERROR 降級為 WARNING
            metadata["error"] = error_msg
            results.append((error_msg, metadata))
        except Exception as e:
            error_msg = f"解析 Visio 文件 {file_path} 時出錯: {str(e)}"
            logger.warning(error_msg)  # 將 ERROR 降級為 WARNING
            metadata["error"] = str(e)
            results.append((error_msg, metadata))
            
        return results


class DocParser(FileParser):
    """DOC 文件解析器"""
    
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        """
        解析 DOC 文件
        
        Args:
            file_path: 文件路徑
            
        Returns:
            包含 (文本, 元數據) 元組的列表
        """
        metadata = self.extract_metadata(file_path)
        text = ""
        
        # 文件大小檢查 - 過小的文件很可能不是有效的DOC
        file_size = os.path.getsize(file_path)
        if file_size < 1024:  # 小於 1KB 的文件可能不是有效的 DOC
            logger.warning(f"文件過小 ({file_size} 字節)，可能不是有效的 DOC 文件: {file_path}")
            metadata["file_status"] = "too_small"
            return [(f"文件過小 ({file_size} 字節)，可能不是有效的 DOC 文件", metadata)]
        
        # 檢查文件頭部以確定是否可能是 DOC 文件
        try:
            with open(file_path, 'rb') as f:
                header = f.read(8)
                if header != b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
                    logger.info(f"文件不是標準 OLE2 格式，將使用替代方法: {file_path}")
                    metadata["file_status"] = "not_ole2"
                    # 對於非 OLE2 文件，直接使用備用方法，跳過 COM
                    extraction_methods = [
                        self._parse_with_libreoffice,
                        self._parse_with_textract,
                        self._parse_with_antiword,
                        self._parse_with_binary_read,
                        self._binary_force_extract
                    ]
                else:
                    metadata["file_status"] = "standard_ole2"
                    # 標準 OLE2 文件使用所有方法
                    extraction_methods = [
                        self._parse_with_com_win32,
                        self._parse_with_libreoffice,
                        self._parse_with_textract,
                        self._parse_with_antiword,
                        self._parse_with_binary_read,
                        self._binary_force_extract
                    ]
        except Exception as e:
            logger.info(f"檢查文件頭部時出錯，假設為非標準格式: {str(e)}")
            metadata["file_status"] = "check_error"
            extraction_methods = [
                self._parse_with_libreoffice,
                self._parse_with_textract,
                self._parse_with_antiword,
                self._parse_with_binary_read,
                self._binary_force_extract
            ]
        
        # 過濾掉不可用的方法，並按成功率排序
        available_methods = []
        for method in extraction_methods:
            method_name = method.__name__
            if parser_stats.is_method_available(method_name):
                available_methods.append((method, parser_stats.get_success_rate(method_name)))
            else:
                logger.debug(f"跳過不可用的解析方法: {method_name}")
        
        # 按成功率降序排序方法
        available_methods.sort(key=lambda x: x[1], reverse=True)
        sorted_methods = [method for method, _ in available_methods]
        
        # 如果沒有可用方法，直接使用二進制讀取
        if not sorted_methods:
            logger.info("沒有可用的解析方法，直接使用二進制讀取")
            sorted_methods = [self._parse_with_binary_read, self._binary_force_extract]
        
        # 記錄嘗試順序
        method_names = [method.__name__ for method in sorted_methods]
        logger.debug(f"解析方法嘗試順序: {', '.join(method_names)}")
        
        errors = []
        for method_index, extraction_method in enumerate(sorted_methods):
            method_name = extraction_method.__name__
            try:
                logger.info(f"嘗試使用方法 {method_index+1}/{len(sorted_methods)} 解析 DOC 文件: {method_name}")
                text = extraction_method(file_path)
                if text and len(text.strip()) > 0:
                    logger.info(f"成功使用 {method_name} 提取文本，共 {len(text)} 個字符")
                    parser_stats.record_success(method_name)
                    break
                else:
                    logger.info(f"使用 {method_name} 提取的文本為空，嘗試下一個方法")
                    parser_stats.record_failure(method_name)
                    errors.append(f"{method_name}: 提取的文本為空")
            except Exception as e:
                error_msg = f"{method_name}: {str(e)}"
                logger.info(f"使用 {method_name} 提取文本失敗: {error_msg}")
                parser_stats.record_failure(method_name)
                errors.append(error_msg)
        
        # 如果所有方法都失敗
        if not text or len(text.strip()) == 0:
            error_summary = "\n".join(errors)
            logger.error(f"使用所有方法解析 DOC 文件失敗 ({file_path})。錯誤摘要:\n{error_summary}")
            
            # 返回收集到的錯誤信息而不是空文本
            metadata["extraction_errors"] = errors
            text = f"無法解析此 DOC 文件。嘗試了 {len(sorted_methods)} 種不同的提取方法，但均失敗。"
        
        return [(text, metadata)]
    
    def _parse_with_com_win32(self, file_path: str) -> str:
        """
        使用 Win32 COM 接口解析 DOC 文件
        
        Args:
            file_path: DOC 文件路徑
            
        Returns:
            提取的文本
        """
        if platform.system() != 'Windows':
            logger.info("非 Windows 系統，無法使用 COM 接口")
            return ""
        
        import pythoncom
        import win32com.client
        from win32com.client import constants
        
        # 確保文件路徑是絕對路徑
        word = None
        doc = None
        temp_path = None
        
        try:
            # 檢查文件格式
            is_valid_ole2 = False
            try:
                # 簡單檢查是否為有效的 OLE2 文件
                with open(file_path, 'rb') as f:
                    header = f.read(8)
                    is_valid_ole2 = header == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'
            
                if not is_valid_ole2:
                    logger.info(f"文件不是標準 OLE2 格式，嘗試其他方法: {file_path}")
                    return ""
            except Exception as check_error:
                logger.info(f"檢查文件格式時出錯: {str(check_error)}")
        
            # 初始化 COM
            pythoncom.CoInitialize()
        
            # 創建 Word 應用實例
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
            
                # 使用數值而常量來避免某些常見錯誤
                try:
                    word.DisplayAlerts = constants.wdAlertsNone
                except:
                    word.DisplayAlerts = 0  # 直接使用數值
            except Exception as word_init_error:
                # 將錯誤從 ERROR 降為 INFO，因為這是預期內的行為
                logger.info(f"無法初始化 Word 應用，將嘗試其他解析方法: {str(word_init_error)}")
                raise
        
            # 打開文檔
            abs_path = os.path.abspath(file_path)
            logger.debug(f"嘗試使用Word COM接口打開: {abs_path}")
        
            # 有些嚴重損壞的文件可能需要先轉換格式
            if not is_valid_ole2:
                # 嘗試創建臨時副本
                import tempfile
                temp_dir = tempfile.mkdtemp()
                temp_path = os.path.join(temp_dir, os.path.basename(file_path))
                try:
                    shutil.copy2(file_path, temp_path)
                    # 使用較寬鬆的打開方式
                    abs_path = temp_path
                except Exception as copy_error:
                    logger.info(f"創建臨時文件失敗: {str(copy_error)}")
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir)
                    temp_path = None
        
            try:
                # 添加額外的打開選項來處理問題文件
                doc = word.Documents.Open(
                    abs_path, 
                    ReadOnly=True,
                    Visible=False,
                    NoEncodingDialog=True,
                    DoNotConvertChartData=True,
                    DoNotAddToRecentFiles=True,
                    ConfirmConversions=False
                )
            except Exception as open_error:
                # 如果標準打開失敗，嘗試使用更復原性的方法
                logger.info(f"標準打開方法失敗: {str(open_error)}，嘗試恢復模式")
                try:
                    # 使用更多參數，嘗試以恢復模式打開
                    doc = word.Documents.Open(
                        abs_path, 
                        ReadOnly=True,
                        Visible=False,
                        NoEncodingDialog=True,
                        DoNotConvertChartData=True,
                        DoNotAddToRecentFiles=True,
                        ConfirmConversions=False,
                        PasswordDocument="",
                        AddToRecentFiles=False,
                        Revert=True,  # 嘗試恢復上次保存的版本
                        RepairMode=True  # 修復模式
                    )
                except Exception as recover_error:
                    logger.error(f"恢復模式打開失敗: {str(recover_error)}")
                    raise
        
            # 提取文本前先禁用不必要的功能以防錯誤
            try:
                # 禁用頁面設置功能可避免一些常見錯誤
                doc.ActiveWindow.View.ReadingLayout = False
            except:
                pass  # 忽略任何設置錯誤
        
            # 提取文本
            text = ""
            try:
                text = doc.Content.Text
            except Exception as e:
                logger.info(f"無法提取文檔內容，嘗試逐段提取: {str(e)}")
                try:
                    # 逐段提取
                    text_parts = []
                    for para in doc.Paragraphs:
                        try:
                            text_parts.append(para.Range.Text)
                        except:
                            pass
                    text = "\n".join(text_parts)
                except Exception as para_error:
                    logger.info(f"逐段提取失敗: {str(para_error)}")
        
            return text
        
        except Exception as e:
            error_msg = str(e)
            if "PageSetup" in error_msg:
                raise Exception(f"Document.PageSetup問題，文件可能不是標準格式: {error_msg}")
            elif "not an OLE2 file" in error_msg or "Failed" in error_msg:
                raise Exception(f"DOC格式問題，可能不是標準OLE2文件: {error_msg}")
            else:
                raise
        
        finally:
            # 清理資源
            try:
                if doc:
                    doc.Close(SaveChanges=False)
            except:
                pass
        
            try:
                if word:
                    word.Quit()
            except:
                pass
        
            # 釋放 COM 資源
            try:
                import pythoncom
                pythoncom.CoUninitialize()
            except:
                pass
        
            # 清理臨時文件
            if temp_path and os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except:
                    pass
    
    def _parse_with_libreoffice(self, file_path: str) -> str:
        """
        使用 LibreOffice 將 DOC 轉換為文本
        
        Args:
            file_path: DOC 文件路徑
            
        Returns:
            轉換後的文本
        """
        import subprocess
        import tempfile
        import random
        import string
        
        # 創建臨時目錄
        temp_dir = tempfile.mkdtemp()
        try:
            # 生成隨機文件名以避免衝突
            random_suffix = ''.join(random.choices(string.ascii_lowercase + string.digits, k=6))
            output_file = os.path.join(temp_dir, f"output_{random_suffix}.txt")
        
            # 檢查是否安裝了 LibreOffice
            libreoffice_paths = [
                r'C:\Program Files\LibreOffice\program\soffice.exe',  # Windows
                r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',  # Windows 32位
                '/usr/bin/libreoffice',  # Linux
                '/usr/bin/soffice',  # Linux alternative
                '/Applications/LibreOffice.app/Contents/MacOS/soffice'  # macOS
            ]
        
            libreoffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path):
                    libreoffice_path = path
                    break
        
            if not libreoffice_path:
                logger.info("找不到 LibreOffice 執行路徑")
                return ""
        
            # 轉換命令
            cmd = [
                libreoffice_path,
                '--headless',
                '--convert-to', 'txt:Text (encoded):UTF8',
                '--outdir', temp_dir,
                file_path
            ]
        
            # 執行轉換
            process = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
            # 檢查是否成功
            if process.returncode != 0:
                err_output = process.stderr.decode('utf-8', errors='ignore')
                raise Exception(f"LibreOffice 轉換失敗: {err_output}")
        
            # 尋找生成的 .txt 文件
            txt_files = [f for f in os.listdir(temp_dir) if f.endswith('.txt')]
            if not txt_files:
                raise Exception("轉換過程未生成文本文件")
        
            # 讀取文本內容
            with open(os.path.join(temp_dir, txt_files[0]), 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        
        finally:
            # 清理臨時文件
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
    
    def _parse_with_textract(self, file_path: str) -> str:
        """
        使用 textract 庫解析 DOC 文件
        
        Args:
            file_path: DOC 文件路徑
            
        Returns:
            文本內容
        """
        try:
            import textract
            text = textract.process(file_path).decode('utf-8', errors='ignore')
            return text
        except ImportError:
            logger.debug("未安裝 textract 庫，無法使用")
            return ""
        except Exception as e:
            logger.info(f"使用 textract 解析 DOC 文件時出錯: {str(e)}")
            raise
    
    def _parse_with_antiword(self, file_path: str) -> str:
        """
        使用 antiword 工具解析 DOC 文件
        
        Args:
            file_path: DOC 文件路徑
            
        Returns:
            文本內容
        """
        try:
            # 檢查 antiword 是否可用
            check_cmd = ["where", "antiword"] if platform.system() == "Windows" else ["which", "antiword"]
            result = subprocess.run(check_cmd, capture_output=True, text=True, check=False)
        
            if result.returncode != 0:
                logger.debug("系統中未找到 antiword 工具")
                return ""
        
            # 使用 antiword 提取文本
            cmd = ["antiword", file_path]
            result = subprocess.run(cmd, capture_output=True, text=True, check=False)
        
            if result.returncode == 0 and result.stdout:
                return result.stdout
            else:
                logger.info(f"antiword 未能提取文本: {result.stderr}")
                return ""
        
        except Exception as e:
            logger.info(f"使用 antiword 提取文本時出錯: {str(e)}")
            raise
    
    def _parse_with_binary_read(self, file_path: str) -> str:
        """
        通過二進制讀取嘗試提取 DOC 文件中的可讀文本
        這是一個最後的求助方法，可能只能提取部分內容
        
        Args:
            file_path: DOC 文件路徑
            
        Returns:
            文本內容
        """
        try:
            # 以二進制模式讀取文件
            with open(file_path, 'rb') as file:
                content = file.read()
        
            # 嘗試以不同編碼解碼內容
            encodings = ['utf-8', 'latin1', 'cp1252', 'gbk', 'big5']
            extracted_text = ""
        
            for encoding in encodings:
                try:
                    # 嘗試解碼
                    decoded_content = content.decode(encoding, errors='ignore')
        
                    # 獲取所有可打印的 ASCII 字符和常見標點符號
                    import re
                    printable_text = re.sub(r'[^\x20-\x7E\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+', 
                                            ' ', decoded_content)
        
                    # 清理連續空格和空行
                    clean_text = re.sub(r'\s+', ' ', printable_text).strip()
        
                    # 如果提取出更多文本，使用這個結果
                    if len(clean_text) > len(extracted_text):
                        extracted_text = clean_text
                except:
                    continue
        
            # 如果提取出足夠長的文本，視為成功
            if len(extracted_text) > 100:
                logger.info(f"通過二進制讀取成功提取了 {len(extracted_text)} 個字符")
                return extracted_text
            else:
                logger.info("二進制讀取未提取出足夠的文本內容")
                return ""
        
        except Exception as e:
            logger.error(f"二進制讀取文件時出錯: {str(e)}")
            raise
    
    def _binary_force_extract(self, file_path: str) -> str:
        """
        強制從二進制 DOC 文件中提取任何可能的文本內容
        這是一個極端的方法，用於處理嚴重損壞的文件
        
        Args:
            file_path: 文件路徑
            
        Returns:
            提取的文本
        """
        try:
            logger.info(f"嘗試從損壞的 DOC 文件強制提取文本: {file_path}")
        
            # 讀取整個文件
            with open(file_path, 'rb') as f:
                content = f.read()
        
            extracted_text = ""
        
            # 1. 嘗試找到包含文本的塊 - Word 文件中通常字符間有 00 字節
            text_blocks = []
            in_text = False
            current_block = b''
        
            for i in range(len(content) - 1):
                if content[i] > 31 and content[i] < 127 and content[i+1] == 0:
                    if not in_text:
                        in_text = True
                        current_block = bytes([content[i]])
                    else:
                        current_block += bytes([content[i]])
                elif in_text and content[i] == 0:
                    continue
                elif in_text:
                    in_text = False
                    if len(current_block) > 10:  # 超過10個字符才認為是有意義的文本
                        text_blocks.append(current_block)
                    current_block = b''
        
            # 合併找到的文本塊
            if text_blocks:
                for block in text_blocks:
                    try:
                        block_text = block.decode('utf-16le', errors='ignore')
                        extracted_text += block_text + "\n"
                    except:
                        pass
        
            # 2. 嘗試不同編碼直接解碼
            if not extracted_text or len(extracted_text.strip()) < 50:
                encodings = ['utf-16le', 'utf-8', 'latin1', 'cp1252', 'gbk', 'big5']
                for encoding in encodings:
                    try:
                        decoded = content.decode(encoding, errors='ignore')
        
                        # 獲取所有可打印的 ASCII 字符和常見標點符號
                        import re
                        printable = re.sub(r'[^\x20-\x7E\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+', 
                                            ' ', decoded)
        
                        # 清理連續空格
                        cleaned = re.sub(r'\s+', ' ', printable).strip()
        
                        # 如果提取出更多可讀文本，使用這個結果
                        if len(cleaned) > len(extracted_text) and len(re.findall(r'[A-Za-z]{3,}', cleaned)) > 10:
                            extracted_text = cleaned
                    except:
                        continue
        
            # 3. 嘗試使用特殊字符序列分隔
            if not extracted_text or len(extracted_text.strip()) < 50:
                # DOC 文件中常見的分隔字符序列
                separators = [b'\x13Faaaa', b'\x13F', b'Microsoft Word', b'MSWordDoc']
                for sep in separators:
                    parts = content.split(sep)
                    if len(parts) > 1:
                        for part in parts[1:]:  # 跳過第一部分
                            try:
                                # 取前 1000 個字節嘗試解碼
                                sample = part[:1000]
                                for encoding in encodings:
                                    try:
                                        decoded = sample.decode(encoding, errors='ignore')
                                        if len(decoded.strip()) > 50 and len(re.findall(r'[A-Za-z]{3,}', decoded)) > 5:
                                            extracted_text += decoded + "\n"
                                    except:
                                        continue
                            except:
                                continue
        
            if not extracted_text or len(extracted_text.strip()) < 30:
                logger.info("強制提取未能獲得足夠的文本內容")
                return ""
        
            logger.info(f"成功從損壞文件中提取了 {len(extracted_text)} 個字符")
            return extracted_text
        
        except Exception as e:
            logger.error(f"強制提取文本時出錯: {str(e)}")
            return ""


class ExcelOldParser(FileParser):
    """舊版 Excel (.xls) 文件解析器"""
    
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        try:
            results = []
            metadata = self.extract_metadata(file_path)
            metadata.update({"document_type": "XLS"})
        
            # 方法1: 嘗試使用 xlrd (注意: 2.0.1+ 不支持 .xls)
            try:
                import xlrd
            
                # 檢查xlrd版本
                xlrd_version = tuple(map(int, xlrd.__version__.split('.')))
                supports_xls = xlrd_version < (2, 0, 0)
            
                if supports_xls:
                    logger.info(f"使用xlrd {xlrd.__version__} 解析XLS文件")
                    try:
                        workbook = xlrd.open_workbook(file_path, formatting_info=False, on_demand=True)
                        sheet_processed = self._process_xlrd_workbook(workbook, metadata)
                        if sheet_processed:
                            return results
                    except Exception as open_error:
                        logger.warning(f"使用xlrd打開XLS文件失敗: {str(open_error)}")
                else:
                    logger.warning(f"當前xlrd版本 {xlrd.__version__} 不支持.xls文件，嘗試其他方法")
            except ImportError:
                logger.warning("未安裝xlrd庫，嘗試其他解析方法")
        
            # 方法2: 嘗試使用 pyexcel (如果已安裝)
            try:
                import pyexcel
                sheet = pyexcel.get_sheet(file_name=file_path)
                if sheet:
                    sheet_text = []
                    for row in sheet:
                        if any(str(cell).strip() for cell in row):
                            row_values = [str(cell) for cell in row]
                            sheet_text.append(" | ".join(row_values))
        
                    if sheet_text:
                        sheet_content = "\n".join(sheet_text)
                        results.append((sheet_content, metadata))
                        return results
            except ImportError:
                logger.warning("未安裝pyexcel庫")
            except Exception as pyexcel_error:
                logger.warning(f"使用pyexcel解析XLS文件失敗: {str(pyexcel_error)}")
        
            # 方法3: Windows平台使用 pywin32
            if platform.system() == "Windows":
                try:
                    import win32com.client
                    excel = None
                    try:
                        excel = win32com.client.Dispatch("Excel.Application")
            
                        # 使用絕對路徑
                        abs_path = os.path.abspath(file_path)
                        logger.info(f"嘗試使用Excel COM接口打開: {abs_path}")
            
                        workbook = excel.Workbooks.Open(abs_path)
                        text_content = []
            
                        for sheet_index in range(1, workbook.Sheets.Count + 1):
                            sheet = workbook.Sheets(sheet_index)
                            used_range = sheet.UsedRange
            
                            # 獲取表格範圍
                            if used_range:
                                for row_index in range(1, min(used_range.Rows.Count, 5000) + 1):
                                    row_values = []
                                    for col_index in range(1, used_range.Columns.Count + 1):
                                        cell_value = str(used_range.Cells(row_index, col_index).Value or "")
                                        row_values.append(cell_value)
            
                                    if any(cell.strip() for cell in row_values):
                                        text_content.append(" | ".join(row_values))
            
                        workbook.Close(SaveChanges=False)
            
                        if results:
                            return results
                    except Exception as com_error:
                        logger.warning(f"Excel COM接口錯誤: {str(com_error)}")
                    finally:
                        # 確保Excel應用程序被關閉
                        if excel:
                            try:
                                excel.Quit()
                            except:
                                pass
                except ImportError:
                    logger.warning("未安裝pywin32，嘗試其他解析方法")
        
            # 方法4: 嘗試使用 LibreOffice 轉換（如果安裝）
            try:
                if shutil.which("soffice"):
                    logger.info("嘗試使用LibreOffice轉換XLS文件")
                    temp_dir = tempfile.mkdtemp()
                    temp_xlsx = os.path.join(temp_dir, "temp.xlsx")
        
                    try:
                        # 轉換為xlsx
                        subprocess.run([
                            "soffice", 
                            "--headless", 
                            "--convert-to", 
                            "xlsx", 
                            "--outdir", 
                            temp_dir, 
                            file_path
                        ], check=True)
        
                        # 如果轉換成功，使用openpyxl讀取
                        if os.path.exists(temp_xlsx):
                            workbook = openpyxl.load_workbook(temp_xlsx, read_only=True, data_only=True)
        
                            for sheet_name in workbook.sheetnames:
                                sheet = workbook[sheet_name]
                                sheet_text = []
            
                                row_count = 0
                                for row in sheet.iter_rows(values_only=True):
                                    if row_count >= 5000:  # 限制行數
                                        sheet_text.append("... (表格過大，僅顯示前5000行)")
                                        break
            
                                    if any(str(cell).strip() if cell is not None else "" for cell in row):
                                        row_values = [str(cell) if cell is not None else "" for cell in row]
                                        sheet_text.append(" | ".join(row_values))
            
                                    row_count += 1
            
                                if sheet_text:
                                    sheet_content = "\n".join(sheet_text)
                                    sheet_metadata = metadata.copy()
                                    sheet_metadata.update({"sheet_name": sheet_name})
                                    results.append((sheet_content, sheet_metadata))
        
                            if results:
                                return results
                    finally:
                        # 清理臨時文件
                        try:
                            shutil.rmtree(temp_dir)
                        except:
                            pass
            except Exception as lo_error:
                logger.warning(f"使用LibreOffice轉換XLS文件失敗: {str(lo_error)}")
        
            # 所有方法都失敗，返回錯誤消息
            logger.error(f"無法解析 .xls 文件 {file_path}，嘗試了所有可行方法")
            results.append(("無法解析此Excel文件，請確保文件未損壞或受保護", metadata))
        
            return results
        
        except Exception as e:
            logger.error(f"解析 XLS 文件 {file_path} 時出錯: {str(e)}")
            import traceback
            logger.error(f"XLS解析詳細錯誤: {traceback.format_exc()}")
        
            # 返回一個簡單的結果而不是空列表
            metadata = self.extract_metadata(file_path)
            metadata.update({"document_type": "XLS", "error": str(e)})
            return [("XLS解析出錯，可能是文件格式問題或權限限制", metadata)]
        
    def _process_xlrd_workbook(self, workbook, metadata):
        """處理xlrd工作簿並將結果添加到results列表中"""
        results = []
        sheet_processed = False
        
        for sheet_index in range(workbook.nsheets):
            try:
                sheet = workbook.sheet_by_index(sheet_index)
                sheet_text = []
            
                # 設置行數限制以防止超大表格
                max_rows = 5000
            
                for row_idx in range(min(sheet.nrows, max_rows)):
                    try:
                        row_values = []
                        for cell in sheet.row_values(row_idx):
                            if cell is None:
                                row_values.append("")
                            elif isinstance(cell, (float, int)) and cell == int(cell):
                                row_values.append(str(int(cell)))
                            else:
                                row_values.append(str(cell))
            
                        if any(filter(bool, row_values)):  # 如果行不為空
                            sheet_text.append(" | ".join(row_values))
                    except Exception as row_error:
                        logger.warning(f"處理XLS行 {row_idx+1} 時出錯: {str(row_error)}")
                        continue
            
                if row_idx >= max_rows - 1 and sheet.nrows > max_rows:
                    sheet_text.append("... (表格過大，僅顯示前5000行)")
            
                if sheet_text:
                    sheet_content = "\n".join(sheet_text)
                    sheet_metadata = metadata.copy()
                    sheet_metadata.update({
                        "sheet_name": sheet.name
                    })
                    results.append((sheet_content, sheet_metadata))
                    sheet_processed = True
            except Exception as sheet_error:
                logger.warning(f"處理XLS表格 {sheet_index} 時出錯: {str(sheet_error)}")
                # 繼續處理下一個表格
        
        return sheet_processed


class PPTParser(FileParser):
    """舊版 PowerPoint (.ppt) 文件解析器"""
    
    def parse(self, file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
        try:
            results = []
            metadata = self.extract_metadata(file_path)
            metadata.update({"document_type": "PPT"})
        
            # 方法1: 嘗試使用 pywin32 (Windows平台)
            if platform.system() == "Windows":
                try:
                    import win32com.client
                    powerpoint = None
                    try:
                        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            
                        # 使用絕對路徑
                        abs_path = os.path.abspath(file_path)
                        logger.info(f"嘗試打開PPT文件: {abs_path}")
            
                        presentation = powerpoint.Presentations.Open(abs_path, False, False, False)
            
                        text_content = []
            
                        for slide_index in range(1, presentation.Slides.Count + 1):
                            slide = presentation.Slides(slide_index)
                            for shape_index in range(1, slide.Shapes.Count + 1):
                                shape = slide.Shapes(shape_index)
                                if shape.HasTextFrame:
                                    if shape.TextFrame.HasText:
                                        text = shape.TextFrame.TextRange.Text.strip()
                                        if text:
                                            text_content.append(text)
            
                        presentation.Close()
            
                        if results:
                            return results
                    except Exception as com_error:
                        logger.warning(f"PowerPoint COM接口錯誤: {str(com_error)}")
                    finally:
                        # 確保PowerPoint應用程序被關閉
                        if powerpoint:
                            try:
                                powerpoint.Quit()
                            except:
                                pass
                except ImportError:
                    logger.warning("未安裝pywin32，嘗試其他解析方法")
        
            # 方法2: 嘗試使用 LibreOffice 轉換（如果安裝）
            try:
                if shutil.which("soffice"):
                    logger.info("嘗試使用LibreOffice轉換PPT文件")
                    temp_dir = tempfile.mkdtemp()
                    temp_pptx = os.path.join(temp_dir, "temp.pptx")
        
                    try:
                        # 轉換為pptx
                        subprocess.run([
                            "soffice", 
                            "--headless", 
                            "--convert-to", 
                            "pptx", 
                            "--outdir", 
                            temp_dir, 
                            file_path
                        ], check=True)
        
                        # 如果轉換成功，使用python-pptx讀取
                        if os.path.exists(temp_pptx):
                            from pptx import Presentation
                            prs = Presentation(temp_pptx)
        
                            text_content = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        text_content.append(shape.text)
        
                            if text_content:
                                results.append(("\n".join(text_content), metadata))
                                return results
                    finally:
                        # 清理臨時文件
                        try:
                            shutil.rmtree(temp_dir)
                        except:
                            pass
            except Exception as lo_error:
                logger.warning(f"使用LibreOffice轉換PPT文件失敗: {str(lo_error)}")
        
            # 方法3: 嘗試使用 textract 庫（如果已安裝）
            try:
                import textract
                text = textract.process(file_path, extension='ppt').decode('utf-8')
                if text.strip():
                    results.append((text, metadata))
                    return results
            except ImportError:
                logger.warning("未安裝textract庫")
            except Exception as textract_error:
                logger.warning(f"使用textract解析PPT文件失敗: {str(textract_error)}")
        
            # 所有方法都失敗，返回錯誤消息
            logger.error(f"無法解析 .ppt 文件 {file_path}，嘗試了所有可行方法")
            results.append(("無法解析此PPT文件，請確保文件未損壞或受保護", metadata))
            return results
        
        except Exception as e:
            logger.error(f"解析 PPT 文件 {file_path} 時出錯: {str(e)}")
            import traceback
            logger.error(f"PPT解析詳細錯誤: {traceback.format_exc()}")
        
            # 返回一個簡單的結果而不是空列表
            metadata = self.extract_metadata(file_path)
            metadata.update({"document_type": "PPT", "error": str(e)})
            return [("PPT文件解析出錯，可能是文件格式問題或權限限制", metadata)]


SUPPORTED_FILE_TYPES=".pdf,.docx,.doc,.xlsx,.xls,.txt,.md,.pptx,.ppt,.csv,.vsdx,.vsd" 